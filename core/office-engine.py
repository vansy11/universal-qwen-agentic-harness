import sys
import json
import os

def read_office(file_path):
    """Reads text from .docx, .xlsx, or .pptx files."""
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    
    if ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        
    elif ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(filename=file_path, data_only=True)
        for sheet in wb.sheetnames:
            text += f"--- Sheet: {sheet} ---\n"
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                text += "\t".join([str(c) if c is not None else "" for c in row]) + "\n"
                
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text += f"--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    else:
        return {"error": "Unsupported file format"}
        
    return {"text": text}

def write_office(file_path, content):
    """Writes text to a .docx, .xlsx, or .pptx file."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".docx":
        from docx import Document
        doc = Document()
        # Split content by newlines to create paragraphs
        for para in content.split("\n"):
            doc.add_paragraph(para)
        doc.save(file_path)
        return {"status": "success", "message": f"Word document saved to {file_path}"}
        
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation()
        # Split content by "---" to create slides
        slides_content = content.split("---")
        for slide_text in slides_content:
            if slide_text.strip():
                slide_layout = prs.slide_layouts[1] # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                shapes = slide.shapes
                shapes.title.text = "Slide Title"
                shapes.placeholders[1].text = slide_text.strip()
        prs.save(file_path)
        return {"status": "success", "message": f"PowerPoint saved to {file_path}"}

    elif ext == ".xlsx":
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        for line in content.split("\n"):
            if not line.strip():
                continue
            cells = line.split("\t") if "\t" in line else line.split(",")
            ws.append([cell.strip() for cell in cells])
        wb.save(file_path)
        return {"status": "success", "message": f"Excel workbook saved to {file_path}"}

    else:
        return {"error": "Unsupported write format"}

if __name__ == "__main__":
    input_data = json.loads(sys.stdin.read())
    func = input_data.get("function")
    file_path = input_data.get("file_path")
    data = input_data.get("data", "")
    
    result = {}
    try:
        if func == "read":
            result = read_office(file_path)
        elif func == "write":
            result = write_office(file_path, data)
        else:
            result = {"error": "Function not supported"}
    except Exception as e:
        result = {"error": str(e)}
        
    print(json.dumps(result))
